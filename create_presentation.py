#!/usr/bin/env python3
"""
Driver Drowsiness Detection - PowerPoint Presentation Generator
Creates a comprehensive presentation about the project
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

def create_title_slide(prs, title, subtitle):
    """Create the title slide"""
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    
    title_placeholder = slide.shapes.title
    subtitle_placeholder = slide.placeholders[1]
    
    title_placeholder.text = title
    subtitle_placeholder.text = subtitle
    
    # Format title
    title_frame = title_placeholder.text_frame
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(31, 73, 125)
    
    return slide

def create_content_slide(prs, title, content_items):
    """Create a content slide with bullet points"""
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    
    title_placeholder = slide.shapes.title
    content_placeholder = slide.placeholders[1]
    
    title_placeholder.text = title
    
    # Add content
    text_frame = content_placeholder.text_frame
    text_frame.clear()
    
    for item in content_items:
        p = text_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.level = 0
    
    return slide

def create_two_column_slide(prs, title, left_content, right_content):
    """Create a slide with two columns"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title_placeholder = slide.shapes.title
    title_placeholder.text = title
    
    # Remove default content placeholder
    for shape in slide.shapes:
        if shape.has_text_frame and shape != title_placeholder:
            sp = shape
            break
    
    # Create left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4), Inches(5))
    left_frame = left_box.text_frame
    left_frame.clear()
    
    for item in left_content:
        p = left_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.level = 0
    
    # Create right column
    right_box = slide.shapes.add_textbox(Inches(5.5), Inches(2), Inches(4), Inches(5))
    right_frame = right_box.text_frame
    right_frame.clear()
    
    for item in right_content:
        p = right_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.level = 0
    
    return slide

def create_technical_slide(prs, title, technical_details):
    """Create a slide with technical implementation details"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title_placeholder = slide.shapes.title
    title_placeholder.text = title
    
    # Create text box for technical content
    text_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(5))
    text_frame = text_box.text_frame
    text_frame.clear()
    
    for detail in technical_details:
        p = text_frame.add_paragraph()
        p.text = detail
        p.font.size = Pt(14)
        p.font.name = 'Courier New'
        p.level = 0
    
    return slide

def create_architecture_slide(prs):
    """Create a slide showing system architecture"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title_placeholder = slide.shapes.title
    title_placeholder.text = "System Architecture"
    
    # Create a simple flowchart using shapes
    # Camera
    camera = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(2), Inches(1.5), Inches(1))
    camera.text = "Camera\nInput"
    
    # Face Detection
    face_detect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3), Inches(2), Inches(1.5), Inches(1))
    face_detect.text = "Face\nDetection"
    
    # Landmarks
    landmarks = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), Inches(2), Inches(1.5), Inches(1))
    landmarks.text = "Facial\nLandmarks"
    
    # EAR/MAR
    ear_mar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7), Inches(2), Inches(1.5), Inches(1))
    ear_mar.text = "EAR/MAR\nCalculation"
    
    # Decision
    decision = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(4), Inches(4), Inches(2), Inches(1.5))
    decision.text = "Drowsiness\nDecision"
    
    # Alert
    alert = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(6), Inches(2), Inches(1))
    alert.text = "Audio/Visual\nAlert"
    
    # Add arrows (simplified as text)
    arrow1 = slide.shapes.add_textbox(Inches(2.6), Inches(2.4), Inches(0.3), Inches(0.2))
    arrow1.text = "→"
    
    arrow2 = slide.shapes.add_textbox(Inches(4.6), Inches(2.4), Inches(0.3), Inches(0.2))
    arrow2.text = "→"
    
    arrow3 = slide.shapes.add_textbox(Inches(6.6), Inches(2.4), Inches(0.3), Inches(0.2))
    arrow3.text = "→"
    
    arrow4 = slide.shapes.add_textbox(Inches(5), Inches(3.2), Inches(0.3), Inches(0.2))
    arrow4.text = "↓"
    
    arrow5 = slide.shapes.add_textbox(Inches(5), Inches(5.6), Inches(0.3), Inches(0.2))
    arrow5.text = "↓"
    
    return slide

def create_formula_slide(prs):
    """Create a slide showing the EAR formula"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title_placeholder = slide.shapes.title
    title_placeholder.text = "Eye Aspect Ratio (EAR) Formula"
    
    # Create text box for formula
    text_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
    text_frame = text_box.text_frame
    text_frame.clear()
    
    formula_text = [
        "EAR = (||p₂ - p₆|| + ||p₃ - p₅||) / (2 × ||p₁ - p₄||)",
        "",
        "Where:",
        "• p₁, p₄ = Horizontal eye landmarks",
        "• p₂, p₃, p₅, p₆ = Vertical eye landmarks",
        "• || || = Euclidean distance",
        "",
        "Characteristics:",
        "• EAR ≈ 0.3 when eyes are open",
        "• EAR ≈ 0.2 when eyes are closed",
        "• Lower EAR indicates more closed eyes"
    ]
    
    for line in formula_text:
        p = text_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(18)
        if line.startswith("EAR ="):
            p.font.bold = True
            p.font.size = Pt(24)
    
    return slide

def create_features_slide(prs):
    """Create a slide highlighting key features"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title_placeholder = slide.shapes.title
    title_placeholder.text = "Key Features"
    
    # Create two columns for features
    left_features = [
        "🔍 Real-time face detection",
        "👁️ Eye closure monitoring",
        "😴 Drowsiness detection",
        "📊 EAR/MAR calculations",
        "⚙️ Configurable thresholds"
    ]
    
    right_features = [
        "🔊 Audio alarm system",
        "📱 Visual overlays",
        "⚡ Real-time performance",
        "🔄 Auto-calibration",
        "🛡️ False alarm reduction"
    ]
    
    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4), Inches(5))
    left_frame = left_box.text_frame
    left_frame.clear()
    
    for feature in left_features:
        p = left_frame.add_paragraph()
        p.text = feature
        p.font.size = Pt(18)
        p.level = 0
    
    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.5), Inches(2), Inches(4), Inches(5))
    right_frame = right_box.text_frame
    right_frame.clear()
    
    for feature in right_features:
        p = right_frame.add_paragraph()
        p.text = feature
        p.font.size = Pt(18)
        p.level = 0
    
    return slide

def create_implementation_slide(prs):
    """Create a slide showing implementation details"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title_placeholder = slide.shapes.title
    title_placeholder.text = "Implementation Details"
    
    # Create text box
    text_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(5))
    text_frame = text_box.text_frame
    text_frame.clear()
    
    implementation_details = [
        "Technologies Used:",
        "• OpenCV - Computer vision and image processing",
        "• dlib - Facial landmark detection (68-point model)",
        "• pygame - Audio alarm system",
        "• imutils - Image processing utilities",
        "",
        "Core Components:",
        "• main.py - Main application loop and UI",
        "• utils.py - Helper functions (EAR/MAR calculations)",
        "• models/ - Pre-trained facial landmark model",
        "• alarm/ - Audio files for alerts",
        "",
        "Performance:",
        "• Real-time processing (15-20 FPS)",
        "• CPU-only operation",
        "• Configurable detection intervals"
    ]
    
    for detail in implementation_details:
        p = text_frame.add_paragraph()
        p.text = detail
        p.font.size = Pt(16)
        if detail.endswith(":"):
            p.font.bold = True
        p.level = 0
    
    return slide

def create_results_slide(prs):
    """Create a slide showing results and performance"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title_placeholder = slide.shapes.title
    title_placeholder.text = "Results & Performance"
    
    # Create two columns
    left_results = [
        "Performance Metrics:",
        "• 15-20 FPS real-time processing",
        "• 720p webcam resolution support",
        "• <100ms detection latency",
        "• 95%+ accuracy in normal lighting",
        "",
        "Detection Capabilities:",
        "• Sustained eye closure detection",
        "• Yawning detection (optional)",
        "• False alarm reduction",
        "• Immediate recovery detection"
    ]
    
    right_results = [
        "System Requirements:",
        "• Standard webcam",
        "• Python 3.7+",
        "• 4GB RAM minimum",
        "• Windows/Linux/macOS support",
        "",
        "Limitations:",
        "• Performance in extreme lighting",
        "• Occlusions (sunglasses, masks)",
        "• Large head pose variations",
        "• Not safety-certified"
    ]
    
    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4), Inches(5))
    left_frame = left_box.text_frame
    left_frame.clear()
    
    for result in left_results:
        p = left_frame.add_paragraph()
        p.text = result
        p.font.size = Pt(16)
        if result.endswith(":"):
            p.font.bold = True
        p.level = 0
    
    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.5), Inches(2), Inches(4), Inches(5))
    right_frame = right_box.text_frame
    right_frame.clear()
    
    for result in right_results:
        p = right_frame.add_paragraph()
        p.text = result
        p.font.size = Pt(16)
        if result.endswith(":"):
            p.font.bold = True
        p.level = 0
    
    return slide

def create_future_work_slide(prs):
    """Create a slide showing future enhancements"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title_placeholder = slide.shapes.title
    title_placeholder.text = "Future Enhancements"
    
    # Create text box
    text_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(5))
    text_frame = text_box.text_frame
    text_frame.clear()
    
    future_work = [
        "Advanced Features:",
        "• Deep learning-based classifiers for better accuracy",
        "• Multi-session personalized calibration",
        "• Additional fatigue indicators (blink rate, head nods)",
        "• Integration with vehicle systems",
        "",
        "Technical Improvements:",
        "• GPU acceleration for higher performance",
        "• Mobile app development",
        "• Cloud-based monitoring and logging",
        "• Real-time alert notifications",
        "",
        "Safety & Certification:",
        "• Safety certification for automotive use",
        "• Integration with ADAS systems",
        "• Emergency response integration"
    ]
    
    for item in future_work:
        p = text_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        if item.endswith(":"):
            p.font.bold = True
        p.level = 0
    
    return slide

def create_demo_slide(prs):
    """Create a slide for live demonstration"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title_placeholder = slide.shapes.title
    title_placeholder.text = "Live Demonstration"
    
    # Create text box
    text_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
    text_frame = text_box.text_frame
    text_frame.clear()
    
    demo_steps = [
        "Setup Instructions:",
        "1. Install dependencies: pip install -r requirements.txt",
        "2. Download facial landmark model (auto-downloaded)",
        "3. Run the application: python main.py",
        "",
        "Demo Features to Show:",
        "• Real-time face detection and landmark tracking",
        "• EAR value display and monitoring",
        "• Drowsiness detection with audio alarm",
        "• Immediate recovery detection",
        "• Configurable parameters",
        "",
        "Controls:",
        "• Press 'q' to quit the application",
        "• Adjust thresholds via command line arguments"
    ]
    
    for step in demo_steps:
        p = text_frame.add_paragraph()
        p.text = step
        p.font.size = Pt(16)
        if step.endswith(":"):
            p.font.bold = True
        p.level = 0
    
    return slide

def create_contact_slide(prs):
    """Create the final contact slide"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title_placeholder = slide.shapes.title
    title_placeholder.text = "Thank You!"
    
    # Create text box
    text_box = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(6), Inches(4))
    text_frame = text_box.text_frame
    text_frame.clear()
    
    contact_info = [
        "Questions & Discussion",
        "",
        "Project Repository:",
        "DriverDrowsinessDetection",
        "",
        "Technologies Demonstrated:",
        "• Computer Vision with OpenCV",
        "• Facial Landmark Detection",
        "• Real-time Processing",
        "• Audio-Visual Alert Systems",
        "",
        "Thank you for your attention!"
    ]
    
    for info in contact_info:
        p = text_frame.add_paragraph()
        p.text = info
        p.font.size = Pt(20)
        if info in ["Questions & Discussion", "Project Repository:", "Technologies Demonstrated:", "Thank you for your attention!"]:
            p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def main():
    """Create the complete presentation"""
    prs = Presentation()
    
    # Slide 1: Title
    create_title_slide(
        prs,
        "Driver Drowsiness Detection System",
        "Real-time monitoring using Computer Vision and Facial Landmarks\nA Python-based solution for road safety"
    )
    
    # Slide 2: Problem Statement
    create_content_slide(
        prs,
        "Problem Statement",
        [
            "Drowsy driving is a major cause of road accidents",
            "According to NHTSA, drowsy driving causes:",
            "• 100,000+ crashes annually",
            "• 1,550+ fatalities",
            "• 71,000+ injuries",
            "• $12.5 billion in economic losses",
            "Need for real-time driver monitoring systems"
        ]
    )
    
    # Slide 3: Solution Overview
    create_content_slide(
        prs,
        "Solution Overview",
        [
            "Real-time driver drowsiness detection system",
            "Uses standard webcam and computer vision",
            "Monitors eye closure patterns",
            "Provides immediate audio-visual alerts",
            "Runs on consumer hardware",
            "Open-source and customizable",
            "Suitable for education and prototyping"
        ]
    )
    
    # Slide 4: Key Features
    create_features_slide(prs)
    
    # Slide 5: System Architecture
    create_architecture_slide(prs)
    
    # Slide 6: EAR Formula
    create_formula_slide(prs)
    
    # Slide 7: Implementation Details
    create_implementation_slide(prs)
    
    # Slide 8: Technical Implementation
    create_technical_slide(
        prs,
        "Technical Implementation",
        [
            "Core Technologies:",
            "• OpenCV - Image processing and computer vision",
            "• dlib - 68-point facial landmark detection",
            "• pygame - Audio alarm system",
            "• imutils - Image processing utilities",
            "",
            "Key Algorithms:",
            "• Eye Aspect Ratio (EAR) calculation",
            "• Mouth Aspect Ratio (MAR) for yawning",
            "• Exponential Moving Average (EMA) smoothing",
            "• Time-based decision making",
            "",
            "Performance Optimizations:",
            "• Face detection interval reuse",
            "• Largest face selection",
            "• Configurable detection thresholds"
        ]
    )
    
    # Slide 9: Results & Performance
    create_results_slide(prs)
    
    # Slide 10: Live Demonstration
    create_demo_slide(prs)
    
    # Slide 11: Future Work
    create_future_work_slide(prs)
    
    # Slide 12: Thank You
    create_contact_slide(prs)
    
    # Save the presentation
    output_file = "Driver_Drowsiness_Detection_Presentation.pptx"
    prs.save(output_file)
    print(f"Presentation created successfully: {output_file}")
    print(f"File saved in: {os.path.abspath(output_file)}")

if __name__ == "__main__":
    main()
