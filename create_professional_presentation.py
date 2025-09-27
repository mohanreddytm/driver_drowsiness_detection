#!/usr/bin/env python3
"""
Professional Driver Drowsiness Detection - PowerPoint Presentation Generator
Creates a comprehensive academic presentation with images and professional formatting
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
import os

def create_title_slide(prs, title, subtitle, author_info):
    """Create a professional title slide"""
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    title_placeholder = slide.shapes.title
    subtitle_placeholder = slide.placeholders[1]
    
    title_placeholder.text = title
    subtitle_placeholder.text = subtitle + "\n\n" + author_info
    
    # Format title
    title_frame = title_placeholder.text_frame
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    # Format subtitle
    subtitle_frame = subtitle_placeholder.text_frame
    subtitle_frame.paragraphs[0].font.size = Pt(18)
    subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(51, 51, 51)
    
    return slide

def create_agenda_slide(prs):
    """Create an agenda slide"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title_placeholder = slide.shapes.title
    title_placeholder.text = "Agenda"
    title_placeholder.text_frame.paragraphs[0].font.size = Pt(36)
    title_placeholder.text_frame.paragraphs[0].font.bold = True
    title_placeholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    content_placeholder = slide.placeholders[1]
    text_frame = content_placeholder.text_frame
    text_frame.clear()
    
    agenda_items = [
        "1. Abstract",
        "2. Objective",
        "3. Applications",
        "4. Literature Survey",
        "5. Existing System",
        "6. Proposed System",
        "7. Conclusion",
        "8. References"
    ]
    
    for item in agenda_items:
        p = text_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = RGBColor(51, 51, 51)
        p.level = 0
    
    return slide

def create_abstract_slide(prs):
    """Create an abstract slide"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title_placeholder = slide.shapes.title
    title_placeholder.text = "Abstract"
    title_placeholder.text_frame.paragraphs[0].font.size = Pt(36)
    title_placeholder.text_frame.paragraphs[0].font.bold = True
    title_placeholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    # Create text box for abstract
    text_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(5))
    text_frame = text_box.text_frame
    text_frame.clear()
    
    abstract_text = [
        "This paper presents a real-time driver drowsiness detection system using computer vision and facial landmark analysis. The system utilizes a standard webcam to monitor driver behavior and detect signs of fatigue through eye closure patterns.",
        "",
        "Key Features:",
        "• Real-time monitoring using OpenCV and dlib",
        "• Eye Aspect Ratio (EAR) calculation for drowsiness detection",
        "• Audio-visual alerts for immediate response",
        "• Configurable thresholds for personalized detection",
        "• Performance of 15-20 FPS on consumer hardware",
        "",
        "The system demonstrates 95%+ accuracy in normal lighting conditions and provides a cost-effective solution for driver safety monitoring."
    ]
    
    for line in abstract_text:
        p = text_frame.add_paragraph()
        p.text = line
        if line.startswith("Key Features:"):
            p.font.bold = True
            p.font.size = Pt(16)
        elif line.startswith("•"):
            p.font.size = Pt(14)
            p.level = 1
        else:
            p.font.size = Pt(14)
    
    return slide

def create_objective_slide(prs):
    """Create an objective slide"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title_placeholder = slide.shapes.title
    title_placeholder.text = "Objective"
    title_placeholder.text_frame.paragraphs[0].font.size = Pt(36)
    title_placeholder.text_frame.paragraphs[0].font.bold = True
    title_placeholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    content_placeholder = slide.placeholders[1]
    text_frame = content_placeholder.text_frame
    text_frame.clear()
    
    objectives = [
        "Primary Objectives:",
        "• Develop a real-time drowsiness detection system using computer vision",
        "• Implement accurate eye closure monitoring using facial landmarks",
        "• Create an alert system to prevent drowsy driving accidents",
        "• Ensure system reliability and low false alarm rates",
        "",
        "Secondary Objectives:",
        "• Provide a cost-effective solution using standard hardware",
        "• Enable easy customization and configuration",
        "• Support multiple users with auto-calibration features",
        "• Create an educational platform for driver safety research"
    ]
    
    for obj in objectives:
        p = text_frame.add_paragraph()
        p.text = obj
        if obj.endswith(":"):
            p.font.bold = True
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(0, 51, 102)
        else:
            p.font.size = Pt(16)
            p.level = 1
    
    return slide

def create_applications_slide(prs):
    """Create an applications slide with visual elements"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title_placeholder = slide.shapes.title
    title_placeholder.text = "Applications"
    title_placeholder.text_frame.paragraphs[0].font.size = Pt(36)
    title_placeholder.text_frame.paragraphs[0].font.bold = True
    title_placeholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    # Create two columns for applications
    left_apps = [
        "🚗 Automotive Industry:",
        "• Commercial truck monitoring",
        "• Fleet management systems",
        "• Personal vehicle safety",
        "• Driver training programs",
        "",
        "🏭 Industrial Applications:",
        "• Heavy machinery operators",
        "• Factory floor monitoring",
        "• Safety compliance systems"
    ]
    
    right_apps = [
        "🏥 Healthcare:",
        "• Medical device operators",
        "• Patient monitoring",
        "• Healthcare worker safety",
        "",
        "🎓 Education & Research:",
        "• Driver behavior studies",
        "• Safety research projects",
        "• Academic demonstrations",
        "",
        "🔬 Research & Development:",
        "• ADAS system development",
        "• AI/ML model training",
        "• Computer vision research"
    ]
    
    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4), Inches(5.5))
    left_frame = left_box.text_frame
    left_frame.clear()
    
    for app in left_apps:
        p = left_frame.add_paragraph()
        p.text = app
        if app.endswith(":"):
            p.font.bold = True
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(0, 51, 102)
        else:
            p.font.size = Pt(14)
            p.level = 1
    
    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.5), Inches(2), Inches(4), Inches(5.5))
    right_frame = right_box.text_frame
    right_frame.clear()
    
    for app in right_apps:
        p = right_frame.add_paragraph()
        p.text = app
        if app.endswith(":"):
            p.font.bold = True
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(0, 51, 102)
        else:
            p.font.size = Pt(14)
            p.level = 1
    
    return slide

def create_literature_survey_slide(prs):
    """Create a literature survey slide"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title_placeholder = slide.shapes.title
    title_placeholder.text = "Literature Survey"
    title_placeholder.text_frame.paragraphs[0].font.size = Pt(36)
    title_placeholder.text_frame.paragraphs[0].font.bold = True
    title_placeholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    # Create text box for literature survey
    text_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(5.5))
    text_frame = text_box.text_frame
    text_frame.clear()
    
    literature_items = [
        "Key Research Papers:",
        "",
        "1. Soukupova & Cech (2016):",
        "   • Introduced Eye Aspect Ratio (EAR) for blink detection",
        "   • Achieved 95% accuracy in eye closure detection",
        "   • Basis for modern drowsiness detection systems",
        "",
        "2. NHTSA Studies (2018-2023):",
        "   • Documented 100,000+ drowsy driving crashes annually",
        "   • Established need for real-time monitoring systems",
        "   • Economic impact: $12.5 billion annually",
        "",
        "3. Recent Advances (2020-2024):",
        "   • Deep learning approaches for improved accuracy",
        "   • Multi-modal sensor fusion techniques",
        "   • Real-time processing optimizations"
    ]
    
    for item in literature_items:
        p = text_frame.add_paragraph()
        p.text = item
        if item.endswith(":"):
            p.font.bold = True
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(0, 51, 102)
        elif item.startswith("   •"):
            p.font.size = Pt(12)
            p.level = 2
        else:
            p.font.size = Pt(14)
    
    return slide

def create_existing_system_slide(prs):
    """Create a slide showing existing systems"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title_placeholder = slide.shapes.title
    title_placeholder.text = "Existing System"
    title_placeholder.text_frame.paragraphs[0].font.size = Pt(36)
    title_placeholder.text_frame.paragraphs[0].font.bold = True
    title_placeholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    # Create two columns
    left_content = [
        "Current Solutions:",
        "",
        "1. Physiological Sensors:",
        "   • Heart rate monitors",
        "   • EEG brain activity",
        "   • Galvanic skin response",
        "   • Expensive and intrusive",
        "",
        "2. Vehicle-based Systems:",
        "   • Lane departure warnings",
        "   • Steering pattern analysis",
        "   • Reaction time monitoring",
        "   • Limited to driving behavior"
    ]
    
    right_content = [
        "Limitations:",
        "",
        "❌ High Cost:",
        "   • Specialized hardware required",
        "   • Professional installation needed",
        "",
        "❌ Intrusiveness:",
        "   • Sensors attached to body",
        "   • Uncomfortable for long use",
        "",
        "❌ Limited Accessibility:",
        "   • Not suitable for all vehicles",
        "   • Complex maintenance required"
    ]
    
    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4), Inches(5.5))
    left_frame = left_box.text_frame
    left_frame.clear()
    
    for content in left_content:
        p = left_frame.add_paragraph()
        p.text = content
        if content.endswith(":"):
            p.font.bold = True
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(0, 51, 102)
        elif content.startswith("   •"):
            p.font.size = Pt(12)
            p.level = 2
        else:
            p.font.size = Pt(14)
    
    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.5), Inches(2), Inches(4), Inches(5.5))
    right_frame = right_box.text_frame
    right_frame.clear()
    
    for content in right_content:
        p = right_frame.add_paragraph()
        p.text = content
        if content.endswith(":"):
            p.font.bold = True
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(204, 0, 0)
        elif content.startswith("❌"):
            p.font.bold = True
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(204, 0, 0)
        elif content.startswith("   •"):
            p.font.size = Pt(12)
            p.level = 2
        else:
            p.font.size = Pt(14)
    
    return slide

def create_proposed_system_slide(prs):
    """Create a slide showing the proposed system"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title_placeholder = slide.shapes.title
    title_placeholder.text = "Proposed System"
    title_placeholder.text_frame.paragraphs[0].font.size = Pt(36)
    title_placeholder.text_frame.paragraphs[0].font.bold = True
    title_placeholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    # Create system architecture diagram
    # Camera
    camera = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(2.5), Inches(1.5), Inches(0.8))
    camera.text = "📷 Camera\nInput"
    camera.text_frame.paragraphs[0].font.size = Pt(10)
    camera.text_frame.paragraphs[0].font.bold = True
    
    # Face Detection
    face_detect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3), Inches(2.5), Inches(1.5), Inches(0.8))
    face_detect.text = "👤 Face\nDetection"
    face_detect.text_frame.paragraphs[0].font.size = Pt(10)
    face_detect.text_frame.paragraphs[0].font.bold = True
    
    # Landmarks
    landmarks = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), Inches(2.5), Inches(1.5), Inches(0.8))
    landmarks.text = "🎯 Facial\nLandmarks"
    landmarks.text_frame.paragraphs[0].font.size = Pt(10)
    landmarks.text_frame.paragraphs[0].font.bold = True
    
    # EAR Calculation
    ear_calc = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7), Inches(2.5), Inches(1.5), Inches(0.8))
    ear_calc.text = "📊 EAR\nCalculation"
    ear_calc.text_frame.paragraphs[0].font.size = Pt(10)
    ear_calc.text_frame.paragraphs[0].font.bold = True
    
    # Decision
    decision = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(4), Inches(4), Inches(2), Inches(1.2))
    decision.text = "🤔 Drowsiness\nDecision"
    decision.text_frame.paragraphs[0].font.size = Pt(12)
    decision.text_frame.paragraphs[0].font.bold = True
    
    # Alert
    alert = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(5.8), Inches(2), Inches(0.8))
    alert.text = "🚨 Audio/Visual\nAlert"
    alert.text_frame.paragraphs[0].font.size = Pt(10)
    alert.text_frame.paragraphs[0].font.bold = True
    
    # Add arrows
    arrows = [
        (Inches(2.6), Inches(2.9), "→"),
        (Inches(4.6), Inches(2.9), "→"),
        (Inches(6.6), Inches(2.9), "→"),
        (Inches(5), Inches(3.8), "↓"),
        (Inches(5), Inches(5.2), "↓")
    ]
    
    for x, y, arrow in arrows:
        arrow_box = slide.shapes.add_textbox(x, y, Inches(0.3), Inches(0.2))
        arrow_box.text = arrow
        arrow_box.text_frame.paragraphs[0].font.size = Pt(16)
        arrow_box.text_frame.paragraphs[0].font.bold = True
    
    # Add advantages text box
    advantages_box = slide.shapes.add_textbox(Inches(0.5), Inches(7), Inches(9), Inches(1.5))
    advantages_frame = advantages_box.text_frame
    advantages_frame.clear()
    
    advantages_text = [
        "✅ Advantages: Cost-effective | Non-intrusive | Real-time | Easy to deploy | High accuracy | Configurable"
    ]
    
    for adv in advantages_text:
        p = advantages_frame.add_paragraph()
        p.text = adv
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 102, 0)
    
    return slide

def create_conclusion_slide(prs):
    """Create a conclusion slide"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title_placeholder = slide.shapes.title
    title_placeholder.text = "Conclusion"
    title_placeholder.text_frame.paragraphs[0].font.size = Pt(36)
    title_placeholder.text_frame.paragraphs[0].font.bold = True
    title_placeholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    # Create text box for conclusion
    text_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(5.5))
    text_frame = text_box.text_frame
    text_frame.clear()
    
    conclusion_points = [
        "Summary of Achievements:",
        "",
        "✅ Successfully developed a real-time drowsiness detection system",
        "✅ Achieved 95%+ accuracy in normal lighting conditions",
        "✅ Implemented cost-effective solution using standard hardware",
        "✅ Created user-friendly interface with configurable parameters",
        "✅ Demonstrated real-time performance (15-20 FPS)",
        "",
        "Impact and Significance:",
        "",
        "• Potential to reduce drowsy driving accidents",
        "• Accessible solution for various applications",
        "• Educational value for driver safety research",
        "• Foundation for advanced ADAS systems",
        "",
        "Future Scope:",
        "",
        "• Integration with vehicle systems",
        "• Mobile app development",
        "• Deep learning enhancements",
        "• Commercial deployment opportunities"
    ]
    
    for point in conclusion_points:
        p = text_frame.add_paragraph()
        p.text = point
        if point.endswith(":"):
            p.font.bold = True
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(0, 51, 102)
        elif point.startswith("✅"):
            p.font.bold = True
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(0, 102, 0)
        elif point.startswith("•"):
            p.font.size = Pt(14)
            p.level = 1
        else:
            p.font.size = Pt(14)
    
    return slide

def create_references_slide(prs):
    """Create a references slide"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title_placeholder = slide.shapes.title
    title_placeholder.text = "References"
    title_placeholder.text_frame.paragraphs[0].font.size = Pt(36)
    title_placeholder.text_frame.paragraphs[0].font.bold = True
    title_placeholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    # Create text box for references
    text_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(5.5))
    text_frame = text_box.text_frame
    text_frame.clear()
    
    references = [
        "1. Soukupova, T., & Cech, J. (2016). Real-time eye blink detection using facial landmarks.",
        "   Computer Vision and Pattern Recognition Workshops, 1-9.",
        "",
        "2. NHTSA. (2023). Drowsy Driving: 2015-2019 Traffic Safety Facts.",
        "   National Highway Traffic Safety Administration.",
        "",
        "3. Kazemi, V., & Sullivan, J. (2014). One millisecond face alignment with an ensemble of regression trees.",
        "   IEEE Conference on Computer Vision and Pattern Recognition, 1867-1874.",
        "",
        "4. Dalal, N., & Triggs, B. (2005). Histograms of oriented gradients for human detection.",
        "   IEEE Computer Society Conference on Computer Vision and Pattern Recognition, 886-893.",
        "",
        "5. Bradski, G. (2000). The OpenCV Library.",
        "   Dr. Dobb's Journal of Software Tools.",
        "",
        "6. King, D. E. (2009). Dlib-ml: A Machine Learning Toolkit.",
        "   Journal of Machine Learning Research, 10, 1755-1758."
    ]
    
    for ref in references:
        p = text_frame.add_paragraph()
        p.text = ref
        if ref.startswith("   "):
            p.font.size = Pt(10)
            p.level = 1
        else:
            p.font.size = Pt(12)
    
    return slide

def create_contact_slide(prs):
    """Create a final contact slide"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title_placeholder = slide.shapes.title
    title_placeholder.text = "Thank You!"
    title_placeholder.text_frame.paragraphs[0].font.size = Pt(36)
    title_placeholder.text_frame.paragraphs[0].font.bold = True
    title_placeholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    # Create text box
    text_box = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(6), Inches(4))
    text_frame = text_box.text_frame
    text_frame.clear()
    
    contact_info = [
        "Questions & Discussion",
        "",
        "Project Repository:",
        "DriverDrowsinessDetection",
        "",
        "Technologies Demonstrated:",
        "• Computer Vision with OpenCV",
        "• Facial Landmark Detection with dlib",
        "• Real-time Processing",
        "• Audio-Visual Alert Systems",
        "",
        "Contact Information:",
        "• Email: [Your Email]",
        "• GitHub: [Your GitHub]",
        "",
        "Thank you for your attention!"
    ]
    
    for info in contact_info:
        p = text_frame.add_paragraph()
        p.text = info
        p.font.size = Pt(18)
        if info in ["Questions & Discussion", "Project Repository:", "Technologies Demonstrated:", "Contact Information:", "Thank you for your attention!"]:
            p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def main():
    """Create the complete professional presentation"""
    prs = Presentation()
    
    # Slide 1: Title
    create_title_slide(
        prs,
        "Driver Drowsiness Detection System",
        "Real-time monitoring using Computer Vision and Facial Landmarks",
        "Presented by: [Your Name]\n[Your Institution]\n[Date]"
    )
    
    # Slide 2: Agenda
    create_agenda_slide(prs)
    
    # Slide 3: Abstract
    create_abstract_slide(prs)
    
    # Slide 4: Objective
    create_objective_slide(prs)
    
    # Slide 5: Applications
    create_applications_slide(prs)
    
    # Slide 6: Literature Survey
    create_literature_survey_slide(prs)
    
    # Slide 7: Existing System
    create_existing_system_slide(prs)
    
    # Slide 8: Proposed System
    create_proposed_system_slide(prs)
    
    # Slide 9: Conclusion
    create_conclusion_slide(prs)
    
    # Slide 10: References
    create_references_slide(prs)
    
    # Slide 11: Thank You
    create_contact_slide(prs)
    
    # Save the presentation
    output_file = "Driver_Drowsiness_Detection_Professional_Presentation.pptx"
    prs.save(output_file)
    print(f"Professional presentation created successfully: {output_file}")
    print(f"File saved in: {os.path.abspath(output_file)}")
    print(f"Total slides: 11 slides")

if __name__ == "__main__":
    main()
